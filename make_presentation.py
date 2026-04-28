"""
make_presentation.py — генерирует presentation.pptx с результатами диплома.
Обновлён: 2026-04-28 МСК

Запуск в Colab:
    !pip install python-pptx -q
    !python /content/drive/MyDrive/diploma/make_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Цвета ────────────────────────────────────────────────────────
C_BG        = RGBColor(0xFF, 0xFF, 0xFF)  # белый
C_ACCENT    = RGBColor(0x1F, 0x47, 0x88)  # тёмно-синий
C_TEXT      = RGBColor(0x1A, 0x1A, 0x1A)  # почти чёрный
C_SUBTEXT   = RGBColor(0x55, 0x55, 0x55)  # серый
C_GREEN     = RGBColor(0x2E, 0x7D, 0x32)  # тёмно-зелёный
C_RED       = RGBColor(0xC6, 0x28, 0x28)  # тёмно-красный
C_HEADER_BG = RGBColor(0x1F, 0x47, 0x88)  # синяя шапка

W = Inches(13.33)  # ширина слайда 16:9
H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # полностью пустой
    return prs.slides.add_slide(layout)


def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=C_TEXT,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def header(slide, title_text):
    """Синяя полоса сверху с заголовком."""
    add_rect(slide, 0, 0, W, Inches(1.1), C_HEADER_BG)
    add_text(slide, title_text,
             Inches(0.4), Inches(0.15), W - Inches(0.8), Inches(0.8),
             size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))


def divider(slide, y):
    add_rect(slide, Inches(0.4), y, W - Inches(0.8), Pt(1.5),
             RGBColor(0xCC, 0xCC, 0xCC))


def add_image_safe(slide, path, x, y, w, h=None):
    """Добавляет изображение если файл существует, иначе рисует placeholder."""
    if os.path.exists(path):
        try:
            if h is None:
                slide.shapes.add_picture(path, x, y, width=w)
            else:
                slide.shapes.add_picture(path, x, y, width=w, height=h)
            return True
        except Exception:
            pass
    # placeholder
    add_rect(slide, x, y, w, h or Inches(2.0), RGBColor(0xDD, 0xEE, 0xFF))
    add_text(slide, os.path.basename(path), x + Inches(0.1), y + Inches(0.1),
             w - Inches(0.2), (h or Inches(2.0)) - Inches(0.2),
             size=11, color=C_SUBTEXT, align=PP_ALIGN.CENTER)
    return False


# ══════════════════════════════════════════════════════════════════
# Слайды
# ══════════════════════════════════════════════════════════════════

prs = new_prs()

# ── 1. Титульный слайд ───────────────────────────────────────────
s = blank_slide(prs)
add_rect(s, 0, 0, W, H, C_ACCENT)

add_text(s, "Система автономного удержания полосы\nна основе нейронных сетей и ПИД-регулятора",
         Inches(0.8), Inches(1.8), W - Inches(1.6), Inches(2.4),
         size=36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

add_text(s, "Нейронная сеть (LaneCNN)  +  ПИД-регулятор",
         Inches(0.8), Inches(3.8), W - Inches(1.6), Inches(0.6),
         size=20, color=RGBColor(0xAA, 0xC4, 0xE8), align=PP_ALIGN.CENTER)

add_rect(s, Inches(3.5), Inches(4.6), Inches(6.3), Pt(1.5),
         RGBColor(0xFF, 0xFF, 0xFF))

add_text(s, "Филиппов Фёдор  |  Группа М8О-404Б  |  2026",
         Inches(0.8), Inches(4.9), W - Inches(1.6), Inches(0.6),
         size=18, color=RGBColor(0xDD, 0xDD, 0xDD), align=PP_ALIGN.CENTER)


# ── 2. Постановка задачи ─────────────────────────────────────────
s = blank_slide(prs)
header(s, "Постановка задачи")

add_text(s, "Цель: автономное удержание полосы движения без GPS и карт — только по изображению с камеры.",
         Inches(0.5), Inches(1.3), W - Inches(1.0), Inches(0.6),
         size=18, color=C_TEXT)

divider(s, Inches(2.05))

# Два блока рядом
for i, (title, lines) in enumerate([
    ("Нейронная сеть (LaneCNN)",
     ["Вход: кадр камеры grayscale 32×100",
      "Выход: ошибка отклонения e ∈ [−1, 1]",
      "Роль: «чёрный ящик» — оценка положения в полосе"]),
    ("ПИД-регулятор",
     ["Вход: ошибка e(t) от нейросети",
      "Выход: управляющий сигнал u(t)",
      "Цель: минимизировать ∫|e(t)| dt",
      "Сглаживает шум, подавляет перерегулирование"]),
]):
    x = Inches(0.5) + i * Inches(6.5)
    add_rect(s, x, Inches(2.2), Inches(6.0), Inches(0.5), C_ACCENT)
    add_text(s, title, x + Inches(0.15), Inches(2.25), Inches(5.7), Inches(0.45),
             size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    for j, line in enumerate(lines):
        add_text(s, f"• {line}", x + Inches(0.2),
                 Inches(2.85) + j * Inches(0.52),
                 Inches(5.7), Inches(0.5), size=15, color=C_TEXT)


# ── 2б. Архитектура агента (pipeline) ────────────────────────────
s = blank_slide(prs)
header(s, "Архитектура агента: как устроен и как обучен")

# Блок-схема пайплайна
pipeline = [
    ("Кадр камеры\n1×32×100", Inches(0.5)),
    ("LaneCNN\n(нейросеть)", Inches(3.0)),
    ("e(t) ∈ [−1, 1]\n(ошибка полосы)", Inches(5.5)),
    ("ПИД-регулятор", Inches(8.0)),
    ("u(t) ∈ [−1, 1]\n(угол руля)", Inches(10.5)),
]
for label, x in pipeline:
    add_rect(s, x, Inches(1.5), Inches(2.2), Inches(1.3), C_ACCENT)
    add_text(s, label, x + Inches(0.1), Inches(1.55),
             Inches(2.0), Inches(1.2), size=13, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    if x < Inches(10.0):
        add_text(s, "→", x + Inches(2.25), Inches(1.95),
                 Inches(0.4), Inches(0.5), size=22, bold=True, color=C_ACCENT)

# Определения под пайплайном
defs = [
    ("e(t)", "Cross-Track Error — отклонение от центра полосы, предсказывается нейросетью LaneCNN. e ∈ [−1, 1]"),
    ("u(t)", "Управляющий сигнал ПИД-регулятора = угол поворота руля. u = Kp·e + Ki·∫e·dt + Kd·de/dt"),
    ("CTE",  "Cross-Track Error integral = (1/T)·Σ|u(t)| — средняя абс. ошибка управления за эпизод (↓ лучше)"),
]
for j, (term, desc) in enumerate(defs):
    y_d = Inches(3.0) + j * Inches(0.52)
    add_rect(s, Inches(0.5), y_d, Inches(1.1), Inches(0.45), C_ACCENT)
    add_text(s, term, Inches(0.55), y_d + Inches(0.04), Inches(1.0), Inches(0.38),
             size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text(s, desc, Inches(1.7), y_d + Inches(0.04), W - Inches(2.2), Inches(0.42),
             size=13, color=C_TEXT)

# Как обучали — две фазы
divider(s, Inches(4.7))
add_text(s, "Как обучали агента:",
         Inches(0.5), Inches(4.78), W - Inches(1.0), Inches(0.38),
         size=15, bold=True, color=C_ACCENT)
phases = [
    ("Фаза 1 — Обучение нейросети (Supervised Learning)",
     "LaneCNN обучается по driving_log.csv: вход — кадр камеры, выход — угол руля из записи вождения человека. "
     "Функция потерь: MSE. Оптимизатор: Adam, lr=3·10⁻⁴. EarlyStopping, patience=10."),
    ("Фаза 2 — Подбор коэффициентов ПИД (Grid Search)",
     "Перебор Kp ∈ [0.1, 2.0], Kd ∈ [0.0, 0.5] по тестовой выборке. "
     "Выбираются коэффициенты с минимальным CTE. Ki = 0.01 фиксирован."),
]
for j, (ph_title, ph_desc) in enumerate(phases):
    y_ph = Inches(5.25) + j * Inches(0.95)
    add_rect(s, Inches(0.5), y_ph, Inches(0.35), Inches(0.35),
             C_GREEN if j == 1 else C_ACCENT)
    add_text(s, str(j + 1), Inches(0.5), y_ph + Inches(0.02), Inches(0.35), Inches(0.32),
             size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text(s, ph_title, Inches(0.95), y_ph, W - Inches(1.45), Inches(0.37),
             size=13, bold=True, color=C_TEXT)
    add_text(s, ph_desc, Inches(0.95), y_ph + Inches(0.38), W - Inches(1.45), Inches(0.52),
             size=12, color=C_SUBTEXT)


# ── 2в. Описание симулятора ───────────────────────────────────────
BASE_DIR = os.path.dirname(os.path.abspath(__file__)) if "__file__" in dir() else os.getcwd()

s = blank_slide(prs)
header(s, "Симулятор Udacity: описание и параметры")

# Левая колонка: описание
add_rect(s, Inches(0.5), Inches(1.25), Inches(6.0), Inches(0.45), C_ACCENT)
add_text(s, "Что такое симулятор",
         Inches(0.65), Inches(1.3), Inches(5.7), Inches(0.4),
         size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
sim_desc = [
    "Udacity Self-Driving Car Simulator (Unity Engine, open-source)",
    "Два режима: запись вождения человека / автономное управление",
    "Трек: jungle — лесная дорога с разметкой полосы",
    "Управление: WebSocket — Python-скрипт отправляет угол руля",
    "Данные записываются в driving_log.csv (кадр → угол руля)",
    "3 камеры: центральная, левая, правая (используем центральную)",
]
for j, line in enumerate(sim_desc):
    add_text(s, f"• {line}", Inches(0.65), Inches(1.85) + j * Inches(0.47),
             Inches(5.7), Inches(0.44), size=13, color=C_TEXT)

# Правая колонка: параметры
add_rect(s, Inches(7.0), Inches(1.25), Inches(5.9), Inches(0.45), C_ACCENT)
add_text(s, "Параметры среды",
         Inches(7.15), Inches(1.3), Inches(5.6), Inches(0.4),
         size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
params = [
    ("Разрешение камеры",      "320 × 160 пикселей"),
    ("Частота кадров",         "30 FPS"),
    ("Вход нейросети",         "grayscale 1 × 32 × 100"),
    ("Нормировка пикселей",    "[0, 255] → [−1, 1]"),
    ("u(t) — угол руля",       "∈ [−1, 1] (непрерывный)"),
    ("Скорость (throttle)",    "0.2 — фиксирована"),
    ("Anti-windup (ПИД)",      "integral_max = 5.0"),
]
for j, (param, val) in enumerate(params):
    y_p = Inches(1.85) + j * Inches(0.47)
    bg_p = RGBColor(0xEE, 0xF4, 0xFF) if j % 2 == 0 else C_BG
    add_rect(s, Inches(7.0), y_p, Inches(5.9), Inches(0.44), bg_p)
    add_text(s, param + ":", Inches(7.1), y_p + Inches(0.05), Inches(2.85), Inches(0.36),
             size=13, bold=True, color=C_ACCENT)
    add_text(s, val, Inches(9.95), y_p + Inches(0.05), Inches(2.8), Inches(0.36),
             size=13, color=C_TEXT)

divider(s, Inches(5.2))
add_text(s, "Изображения из записи вождения (ноутбук 01 — анализ датасета):",
         Inches(0.5), Inches(5.28), W - Inches(1.0), Inches(0.36),
         size=13, bold=True, color=C_ACCENT)
img1 = os.path.join(BASE_DIR, "результаты ноутбука 01.png")
add_image_safe(s, img1, Inches(0.5), Inches(5.7), W - Inches(1.0), Inches(1.5))


# ── 3. Данные и модель ───────────────────────────────────────────
s = blank_slide(prs)
header(s, "Данные и архитектура модели")

# Датасет
add_rect(s, Inches(0.5), Inches(1.25), Inches(5.8), Inches(0.45), C_ACCENT)
add_text(s, "Датасет Udacity (jungle)", Inches(0.65), Inches(1.3),
         Inches(5.5), Inches(0.4), size=16, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF))
for j, line in enumerate([
    "3 403 записей  →  1 753 сэмпла (стратификация по углу руля)",
    "16.8% «в полосе» (|e| < 0.15)  |  83.2% «вне полосы»",
    "Среда записи: Windows → абсолютные пути исправлены автоматически",
]):
    add_text(s, f"• {line}", Inches(0.65), Inches(1.85) + j * Inches(0.48),
             Inches(5.5), Inches(0.45), size=14, color=C_TEXT)

# Архитектура
add_rect(s, Inches(6.8), Inches(1.25), Inches(6.0), Inches(0.45), C_ACCENT)
add_text(s, "Архитектура LaneCNN", Inches(6.95), Inches(1.3),
         Inches(5.7), Inches(0.4), size=16, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF))
arch_lines = [
    "Вход: 1×32×100 (grayscale, нормировка [-1,1])",
    "Conv(8) → ReLU → MaxPool  ×3",
    "Flatten → Dense(256) → Dropout(0.3)",
    "Dense(64) → Dense(1) → Tanh",
    "Выход: угол руля ∈ [−1, 1]",
]
for j, line in enumerate(arch_lines):
    add_text(s, f"• {line}", Inches(6.95), Inches(1.85) + j * Inches(0.48),
             Inches(5.7), Inches(0.45), size=14, color=C_TEXT)

divider(s, Inches(4.15))
add_text(s, "Обучение: Adam (lr=3×10⁻⁴) · MSE Loss · EarlyStopping · ReduceLROnPlateau · Google Colab GPU T4",
         Inches(0.5), Inches(4.22), W - Inches(1.0), Inches(0.4),
         size=13, color=C_SUBTEXT)
add_text(s, "Распределение сэмплов по углам руля (до / после балансировки):",
         Inches(0.5), Inches(4.72), W - Inches(1.0), Inches(0.35),
         size=13, bold=True, color=C_ACCENT)
add_image_safe(s, os.path.join(BASE_DIR, "outputs", "angle_distribution.png"),
               Inches(0.5), Inches(5.12), W - Inches(1.0), Inches(2.1))


# ── 3б. Функция потерь (MSE) + ПИД ──────────────────────────────
s = blank_slide(prs)
header(s, "Функция потерь и ПИД-регулятор")

# ── Левая колонка: MSE ──
add_rect(s, Inches(0.5), Inches(1.25), Inches(6.0), Inches(0.45), C_ACCENT)
add_text(s, "Функция потерь: MSE",
         Inches(0.65), Inches(1.3), Inches(5.7), Inches(0.4),
         size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

add_rect(s, Inches(0.6), Inches(1.85), Inches(5.8), Inches(0.75),
         RGBColor(0xE8, 0xF0, 0xFE))
add_text(s, "L(θ) = (1/N) · Σ ( ŷᵢ − yᵢ )²",
         Inches(0.7), Inches(1.9), Inches(5.6), Inches(0.65),
         size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

for j, line in enumerate([
    "ŷᵢ = f(xᵢ; θ) — предсказание нейросети (LaneCNN)",
    "yᵢ  — реальный угол руля из датасета ∈ [−1, 1]",
    "N   — размер батча (batch_size = 32)",
    "Оптимизатор: Adam, lr = 3·10⁻⁴, weight_decay = 10⁻⁴",
]):
    add_text(s, f"• {line}", Inches(0.65), Inches(2.75) + j * Inches(0.46),
             Inches(5.8), Inches(0.44), size=13, color=C_TEXT)

# Примеры значений функции потерь
add_text(s, "Примеры значений функции потерь:",
         Inches(0.65), Inches(4.65), Inches(5.8), Inches(0.38),
         size=13, bold=True, color=C_TEXT)
example_rows = [
    ("Кадр", "ŷ (предсказ.)", "y (реальный)", "(ŷ−y)²"),
    ("#1 — прямая", " 0.03", " 0.00", "0.0009"),
    ("#2 — поворот", "−0.45", "−0.52", "0.0049"),
    ("#3 — выезд",  " 0.61", " 0.38", "0.0529"),
    ("Среднее MSE",  "—",    "—",     "0.0196"),
]
col_xs2 = [Inches(0.65), Inches(2.65), Inches(4.05), Inches(5.5)]
col_ws2 = [Inches(1.9),  Inches(1.3),  Inches(1.35), Inches(0.9)]
for ri, row in enumerate(example_rows):
    y2 = Inches(5.1) + ri * Inches(0.38)
    bg2 = RGBColor(0xE8, 0xEE, 0xF7) if ri == 0 else (
          RGBColor(0xE8, 0xF5, 0xE9) if ri == 4 else
          (RGBColor(0xF5, 0xF5, 0xF5) if ri % 2 == 0 else C_BG))
    add_rect(s, Inches(0.6), y2, Inches(5.8), Inches(0.36), bg2)
    for ci, (val, cx, cw) in enumerate(zip(row, col_xs2, col_ws2)):
        add_text(s, val, cx, y2 + Inches(0.02), cw, Inches(0.34),
                 size=12, bold=(ri == 0 or ri == 4), color=C_ACCENT if ri == 0 else C_TEXT)

# ── Правая колонка: ПИД ──
add_rect(s, Inches(7.0), Inches(1.25), Inches(5.9), Inches(0.45), C_ACCENT)
add_text(s, "ПИД-регулятор: звенья и параметры",
         Inches(7.15), Inches(1.3), Inches(5.6), Inches(0.4),
         size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

add_rect(s, Inches(7.1), Inches(1.85), Inches(5.7), Inches(0.75),
         RGBColor(0xE8, 0xF0, 0xFE))
add_text(s, "u(t) = Kp·e(t) + Ki·∫e·dt + Kd·de/dt",
         Inches(7.15), Inches(1.9), Inches(5.6), Inches(0.65),
         size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Таблица звеньев
zvenia_headers = ["Звено", "Формула", "Роль", "Значение"]
zvenia_rows = [
    ("P — пропорц.", "Kp · e(t)",         "Быстрая реакция",          "Kp = 0.8"),
    ("I — интегр.",  "Ki · Σe·Δt",        "Убирает статич. ошибку",   "Ki = 0.01"),
    ("D — диффер.", "Kd · Δe/Δt",        "Демпфирует колебания",     "Kd = 0.15"),
]
z_col_x = [Inches(7.1),  Inches(8.6),  Inches(9.95), Inches(11.5)]
z_col_w = [Inches(1.4),  Inches(1.25), Inches(1.45), Inches(1.2)]
add_rect(s, Inches(7.1), Inches(2.75), Inches(5.7), Inches(0.38),
         RGBColor(0xE8, 0xEE, 0xF7))
for ci, (hdr, cx, cw) in enumerate(zip(zvenia_headers, z_col_x, z_col_w)):
    add_text(s, hdr, cx, Inches(2.77), cw, Inches(0.35),
             size=12, bold=True, color=C_ACCENT)
for ri, row in enumerate(zvenia_rows):
    y3 = Inches(3.2) + ri * Inches(0.5)
    bg3 = RGBColor(0xF5, 0xF5, 0xF5) if ri % 2 == 0 else C_BG
    add_rect(s, Inches(7.1), y3, Inches(5.7), Inches(0.48), bg3)
    for ci, (val, cx, cw) in enumerate(zip(row, z_col_x, z_col_w)):
        add_text(s, val, cx, y3 + Inches(0.05), cw, Inches(0.42),
                 size=12, color=C_TEXT)

for j, line in enumerate([
    "Δt = 1/30 с (30 FPS)  |  anti-windup: integral_max = 5.0",
    "CTE = (1/T)·Σ|u(t)| — цель минимизировать",
    "Подбор Kp, Ki, Kd — grid search по минимуму CTE",
]):
    add_text(s, f"• {line}", Inches(7.1), Inches(4.9) + j * Inches(0.43),
             Inches(5.7), Inches(0.41), size=13, color=C_TEXT)


# ── 3в. Кривые обучения (графики из ноутбуков) ───────────────────
s = blank_slide(prs)
header(s, "Кривые обучения и результаты предобработки")

add_text(s, "Функция потерь MSE на train/val снижается с каждой эпохой. "
            "Ниже — результаты из ноутбуков 02 (предобработка) и 03 (обучение).",
         Inches(0.5), Inches(1.25), W - Inches(1.0), Inches(0.45),
         size=14, color=C_SUBTEXT)

img2 = os.path.join(BASE_DIR, "outputs", "training_curves.png")
img3 = os.path.join(BASE_DIR, "результаты ноутбука 03.png")

add_text(s, "Кривые обучения Adam — train/val MSE Loss по эпохам:",
         Inches(0.5), Inches(1.82), Inches(6.0), Inches(0.35),
         size=13, bold=True, color=C_ACCENT)
add_image_safe(s, img2, Inches(0.5), Inches(2.2), Inches(6.0), Inches(4.8))

add_text(s, "Ноутбук 03 — кривые loss (MSE train/val по эпохам):",
         Inches(6.9), Inches(1.82), Inches(6.0), Inches(0.35),
         size=13, bold=True, color=C_ACCENT)
add_image_safe(s, img3, Inches(6.9), Inches(2.2), Inches(6.0), Inches(4.8))


# ── 8. Сравнение оптимизаторов ------------------------------------
s = blank_slide(prs)
header(s, "Сравнение оптимизаторов: Adam (baseline) vs SGD vs RMSprop vs Adagrad")

add_text(s, "10 эпох · lr = 3×10⁻⁴ · seed = 42 · Adam — базовый, выделен жирной линией",
         Inches(0.5), Inches(1.25), W - Inches(1.0), Inches(0.4),
         size=14, color=C_SUBTEXT)

add_image_safe(s, os.path.join(BASE_DIR, "outputs", "optimizer_comparison.png"),
               Inches(0.5), Inches(1.75), W - Inches(1.0), Inches(4.3))

divider(s, Inches(6.15))
add_text(s,
    "• Adam сходится быстрее и стабильнее на задаче регрессии угла руля\n"
    "• SGD требует тонкой настройки lr/momentum, медленнее стартует\n"
    "• RMSprop и Adagrad уступают Adam по итоговому val loss",
    Inches(0.5), Inches(6.25), W - Inches(1.0), Inches(1.0),
    size=14, color=C_TEXT)


# ── 9. CrossEntropy vs MSE ----------------------------------------
s = blank_slide(prs)
header(s, "CrossEntropy (BCE) vs MSE: сравнение функций потерь")

add_text(s, "LaneCNNClassifier (Sigmoid) + BCELoss  vs  LaneCNN (Tanh) + MSELoss · 10 эпох · Adam",
         Inches(0.5), Inches(1.25), W - Inches(1.0), Inches(0.4),
         size=14, color=C_SUBTEXT)

add_image_safe(s, os.path.join(BASE_DIR, "outputs", "crossentropy_vs_mse.png"),
               Inches(0.5), Inches(1.75), Inches(7.8), Inches(3.8))

add_rect(s, Inches(8.6), Inches(1.75), Inches(4.3), Inches(0.45), C_ACCENT)
add_text(s, "Ключевые различия", Inches(8.75), Inches(1.8), Inches(4.0), Inches(0.4),
         size=15, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

diff_lines = [
    ("MSELoss",  "Регрессия e в [-1,1]\nВыход: Tanh → ПИД-вход"),
    ("BCELoss",  "Бинарная классификация\nВыход: Sigmoid → вероятность"),
    ("Вывод",   "BCE сопоставима по accuracy\nMSE удобнее для ПИД"),
]
for j, (term, desc) in enumerate(diff_lines):
    y_d = Inches(2.35) + j * Inches(1.15)
    add_rect(s, Inches(8.6), y_d, Inches(1.1), Inches(0.42), C_ACCENT)
    add_text(s, term, Inches(8.62), y_d + Inches(0.04), Inches(1.05), Inches(0.36),
             size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text(s, desc, Inches(9.8), y_d, Inches(3.0), Inches(1.05), size=13, color=C_TEXT)

add_text(s, "В агент (ПИД) по-прежнему подаётся LaneCNN (Tanh) с непрерывным e в [-1, 1].",
         Inches(0.5), Inches(6.2), W - Inches(1.0), Inches(0.55),
         size=13, bold=True, color=C_ACCENT)


# ── 10. Результаты ────────────────────────────────────────────────
s = blank_slide(prs)
header(s, "Результаты: сравнение методов")

add_text(s, "Тестовый трек: 263 кадра (15% от выборки)  |  Метрика CTE = mean(|u(t)|)",
         Inches(0.5), Inches(1.2), W - Inches(1.0), Inches(0.45),
         size=15, color=C_SUBTEXT)

# Таблица
col_x = [Inches(0.5), Inches(5.2), Inches(7.8), Inches(10.4)]
col_w = [Inches(4.5), Inches(2.4), Inches(2.4), Inches(2.3)]
headers_t = ["Метод", "CTE ↓", "В полосе ↑", "Награда"]
rows = [
    ("Только нейросеть",           "0.3258", "20.2%",  "−156",  C_RED),
    ("Нейросеть + ПИД (базовые)",  "0.8359", "2.7%",   "−248",  C_RED),
    ("Нейросеть + ПИД (оптим.)",   "0.0326", "100.0%", "+262",  C_GREEN),
]

# Заголовок таблицы
add_rect(s, Inches(0.5), Inches(1.8), W - Inches(1.0), Inches(0.5),
         RGBColor(0xE8, 0xEE, 0xF7))
for i, (hdr, cx, cw) in enumerate(zip(headers_t, col_x, col_w)):
    add_text(s, hdr, cx + Inches(0.1), Inches(1.82), cw, Inches(0.45),
             size=15, bold=True, color=C_ACCENT)

# Строки
for row_i, (method, cte, pct, rew, color) in enumerate(rows):
    y = Inches(2.4) + row_i * Inches(0.72)
    bg = RGBColor(0xF5, 0xF5, 0xF5) if row_i % 2 == 0 else C_BG
    add_rect(s, Inches(0.5), y, W - Inches(1.0), Inches(0.68), bg)
    vals = [method, cte, pct, rew]
    for i, (val, cx, cw) in enumerate(zip(vals, col_x, col_w)):
        c = color if i > 0 and row_i == 2 else (C_RED if i > 0 and row_i < 2 else C_TEXT)
        b = (row_i == 2)
        add_text(s, val, cx + Inches(0.1), y + Inches(0.1), cw, Inches(0.5),
                 size=15, bold=b, color=c)

add_text(s, "★  Оптимальный ПИД снижает CTE в 10× и обеспечивает 100% удержание полосы",
         Inches(0.5), Inches(4.6), W - Inches(1.0), Inches(0.5),
         size=16, bold=True, color=C_GREEN)


# ── 5. Grid Search и выводы ──────────────────────────────────────
s = blank_slide(prs)
header(s, "Grid Search vs Bayesian Optimization (Optuna)")

# Grid Search
add_rect(s, Inches(0.5), Inches(1.25), Inches(5.8), Inches(0.45), C_ACCENT)
add_text(s, "Grid Search: Kp ∈ [0.1, 2.0], Kd ∈ [0.0, 0.5]",
         Inches(0.65), Inches(1.3), Inches(5.5), Inches(0.4),
         size=15, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
gs_lines = [
    "Оптимум: малый Kp (0.1–0.4), Kd ≈ 0",
    "Большой Kd усиливает шум нейросети → CTE растёт",
    "Большой Kp вызывает перерегулирование",
    "ПИД работает как сглаживатель, а не активный регулятор",
]
for j, line in enumerate(gs_lines):
    add_text(s, f"• {line}", Inches(0.65), Inches(1.85) + j * Inches(0.5),
             Inches(5.5), Inches(0.48), size=14, color=C_TEXT)

# Optuna
add_rect(s, Inches(6.8), Inches(1.25), Inches(6.0), Inches(0.45), C_ACCENT)
add_text(s, "Bayesian Optimization (Optuna, 50 попыток)",
         Inches(6.95), Inches(1.3), Inches(5.7), Inches(0.4),
         size=15, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

add_image_safe(s, os.path.join(BASE_DIR, "outputs", "optuna_convergence.png"),
               Inches(6.8), Inches(1.85), Inches(6.0), Inches(3.3))

add_text(s, "Grid Search vs Optuna:", Inches(6.8), Inches(5.28), Inches(6.0), Inches(0.35),
         size=13, bold=True, color=C_ACCENT)
for j, line in enumerate([
    "Grid Search: 72 комбинации (Kp×Kd), слепой перебор",
    "Optuna: 50 попыток, суррогатная байесовская модель",
    "Optuna находит оптимум быстрее при равном CTE",
]):
    add_text(s, f"• {line}", Inches(6.8), Inches(5.72) + j * Inches(0.44),
             Inches(6.0), Inches(0.42), size=13, color=C_TEXT)


# ── 12. Выводы ----------------------------------------------------
s = blank_slide(prs)
header(s, "Выводы")

concl_12 = [
    ("•", "Нейросеть без ПИД: нестабильный сигнал, 20% в полосе",           False),
    ("•", "Неверные коэффициенты ПИД ухудшают результат (CTE x2.5)",         False),
    ("★", "Оптимальный ПИД: CTE / 10, 100% удержание полосы",               True),
    ("•", "Adam сходится быстрее конкурентов на задаче регрессии угла",       False),
    ("•", "BCE сопоставима с MSE по accuracy, MSE удобнее для ПИД-входа",    False),
    ("•", "Optuna находит оптимум за 50 попыток против 72 у Grid Search",     False),
    ("•", "Подход применим как базовый модуль реальной ADAS",                 False),
]
for j, (sym, line, star) in enumerate(concl_12):
    add_text(s, f"{sym}  {line}",
             Inches(0.8), Inches(1.4) + j * Inches(0.66),
             W - Inches(1.6), Inches(0.62), size=16,
             bold=star, color=C_GREEN if star else C_TEXT)


# ── Сохранение + автоскачивание ──────────────────────────────────
import glob as _glob

def _find_diploma_dir():
    if "__file__" in dir():
        return os.path.dirname(os.path.abspath(__file__))
    # Ищем папку diploma на Google Drive (любое имя аккаунта)
    for pattern in [
        "/content/drive/MyDrive/diploma",
        "/content/drive/My Drive/diploma",
        "/content/drive/*/diploma",
        "/content/drive/*/*/diploma",
    ]:
        found = _glob.glob(pattern)
        if found:
            return found[0]
    return "/content"  # fallback — корень Colab

OUT = os.path.join(_find_diploma_dir(), "presentation.pptx")
prs.save(OUT)
print(f"Сохранено: {OUT}")

# Автоскачивание в браузер (работает только в Colab)
try:
    from google.colab import files
    files.download(OUT)
    print("Скачивание запущено ↓")
except Exception:
    print("(автоскачивание недоступно — запустите вручную)")
